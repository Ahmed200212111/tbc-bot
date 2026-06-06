import logging
import os
import copy
from telegram import Update, ReplyKeyboardMarkup, ReplyKeyboardRemove
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    MessageHandler,
    ConversationHandler,
    ContextTypes,
    filters,
)
from pptx import Presentation
from pptx.util import Pt
from datetime import datetime

# ─── إعدادات ────────────────────────────────────────────────────────────────
TOKEN = "8990357714:AAGgbLbTHKTi2rQ5FkcKTTJ8pO1d41rb4Kw"
TEMPLATE_PATH = "template.pptx"

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

# ─── مراحل المحادثة ──────────────────────────────────────────────────────────
TITLE, SUBTITLE, POINTS, CONFIRM = range(4)


# ─── دوال المساعدة ───────────────────────────────────────────────────────────

def replace_text_in_shape(shape, replacements: dict):
    """استبدال النصوص داخل شكل مع الحفاظ على التنسيق."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for placeholder, value in replacements.items():
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)


def generate_pptx(data: dict) -> str:
    """إنشاء ملف PPTX من القالب وإعادة مسار الملف."""
    prs = Presentation(TEMPLATE_PATH)

    title_text    = data.get("title", "عنوان التقرير")
    subtitle_text = data.get("subtitle", "")
    points        = data.get("points", [])
    date_text     = datetime.now().strftime("%Y-%m-%d")

    replacements = {
        "{{TITLE}}":    title_text,
        "{{SUBTITLE}}": subtitle_text,
        "{{DATE}}":     date_text,
    }

    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)

            # ملء نقاط المحتوى إذا وُجد placeholder مخصص
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "{{POINTS}}" in run.text:
                            run.text = run.text.replace(
                                "{{POINTS}}", "\n".join(f"• {p}" for p in points)
                            )

    # إذا كانت هناك شرائح إضافية مطلوبة لكل نقطة
    if len(prs.slides) > 1 and points:
        template_slide = prs.slides[1]
        slide_layout   = template_slide.slide_layout

        for i, point in enumerate(points[1:], start=2):
            new_slide = prs.slides.add_slide(slide_layout)
            for shape in new_slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    shape.text = point

    output_path = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output_path)
    return output_path


# ─── معالجات المحادثة ────────────────────────────────────────────────────────

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "👋 مرحباً! أنا بوت توليد تقارير PowerPoint.\n\n"
        "أرسل /report لبدء إنشاء تقرير جديد.\n"
        "أرسل /cancel في أي وقت للإلغاء."
    )


async def report_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data.clear()
    await update.message.reply_text(
        "📝 *إنشاء تقرير جديد*\n\nما هو *عنوان* التقرير؟",
        parse_mode="Markdown",
    )
    return TITLE


async def get_title(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["title"] = update.message.text.strip()
    await update.message.reply_text(
        "✏️ ما هو *العنوان الفرعي* (أو أرسل `-` لتخطيه)؟",
        parse_mode="Markdown",
    )
    return SUBTITLE


async def get_subtitle(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = update.message.text.strip()
    context.user_data["subtitle"] = "" if text == "-" else text
    context.user_data["points"] = []
    await update.message.reply_text(
        "📌 أرسل *نقاط المحتوى* واحدة تلو الأخرى.\n"
        "عند الانتهاء أرسل /done",
        parse_mode="Markdown",
    )
    return POINTS


async def get_points(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["points"].append(update.message.text.strip())
    count = len(context.user_data["points"])
    await update.message.reply_text(f"✅ تمت إضافة النقطة {count}. أرسل نقطة أخرى أو /done للمتابعة.")
    return POINTS


async def points_done(update: Update, context: ContextTypes.DEFAULT_TYPE):
    data = context.user_data
    if not data.get("points"):
        await update.message.reply_text("⚠️ يرجى إضافة نقطة واحدة على الأقل.")
        return POINTS

    summary = (
        f"📋 *ملخص التقرير:*\n\n"
        f"• العنوان: {data['title']}\n"
        f"• العنوان الفرعي: {data['subtitle'] or '—'}\n"
        f"• عدد النقاط: {len(data['points'])}\n\n"
        f"هل تريد إنشاء الملف؟"
    )
    keyboard = ReplyKeyboardMarkup([["✅ نعم", "❌ لا"]], resize_keyboard=True, one_time_keyboard=True)
    await update.message.reply_text(summary, parse_mode="Markdown", reply_markup=keyboard)
    return CONFIRM


async def confirm(update: Update, context: ContextTypes.DEFAULT_TYPE):
    choice = update.message.text.strip()
    if "نعم" in choice or "yes" in choice.lower():
        await update.message.reply_text("⏳ جارٍ إنشاء الملف...", reply_markup=ReplyKeyboardRemove())
        try:
            output_path = generate_pptx(context.user_data)
            with open(output_path, "rb") as f:
                await update.message.reply_document(
                    document=f,
                    filename=os.path.basename(output_path),
                    caption="✅ تم إنشاء التقرير بنجاح!",
                )
            os.remove(output_path)
        except Exception as e:
            logger.error(f"Error generating PPTX: {e}")
            await update.message.reply_text(f"❌ حدث خطأ أثناء إنشاء الملف:\n{e}")
    else:
        await update.message.reply_text("تم الإلغاء. أرسل /report لبدء من جديد.", reply_markup=ReplyKeyboardRemove())
    return ConversationHandler.END


async def cancel(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data.clear()
    await update.message.reply_text("❌ تم إلغاء العملية.", reply_markup=ReplyKeyboardRemove())
    return ConversationHandler.END


# ─── نقطة الدخول ─────────────────────────────────────────────────────────────

def main():
    app = ApplicationBuilder().token(TOKEN).build()

    conv_handler = ConversationHandler(
        entry_points=[CommandHandler("report", report_start)],
        states={
            TITLE:   [MessageHandler(filters.TEXT & ~filters.COMMAND, get_title)],
            SUBTITLE:[MessageHandler(filters.TEXT & ~filters.COMMAND, get_subtitle)],
            POINTS:  [
                MessageHandler(filters.TEXT & ~filters.COMMAND, get_points),
                CommandHandler("done", points_done),
            ],
            CONFIRM: [MessageHandler(filters.TEXT & ~filters.COMMAND, confirm)],
        },
        fallbacks=[CommandHandler("cancel", cancel)],
    )

    app.add_handler(CommandHandler("start", start))
    app.add_handler(conv_handler)

    logger.info("🤖 البوت يعمل...")
    app.run_polling()


if __name__ == "__main__":
    main()
