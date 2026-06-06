"""
بوت تلغرام لتوليد تقارير PPTX الفنية
يعتمد على نموذج من Google Drive ويسأل المستخدم عن بيانات المدرسة والزيارة
"""

import os
import io
import logging
import requests
import tempfile
from datetime import datetime
from typing import Optional

from telegram import Update, ReplyKeyboardMarkup, ReplyKeyboardRemove
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    ConversationHandler,
    ContextTypes,
    filters,
)
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from PIL import Image

# ─── إعدادات ────────────────────────────────────────────────────────────────

TOKEN = "8990357714:AAGgbLbTHKTi2rQ5FkcKTTJ8pO1d41rb4Kw"
GDRIVE_FILE_ID = "1McwIYDlaRY9dxE82f3d_d0VNC2DuN9bL"
TEMPLATE_PATH = "template.pptx"

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

# ─── مراحل المحادثة ──────────────────────────────────────────────────────────

(
    SCHOOL_NAME,
    MINISTRY_NUM,
    VISIT_DATE,
    NOTES,
    ACTIONS,
    RECOMMENDATIONS,
    PHOTOS,
    CONFIRM,
) = range(8)

# ─── تحميل النموذج من Google Drive ──────────────────────────────────────────

def download_template() -> bool:
    """تحميل ملف النموذج من Google Drive عند التشغيل"""
    url = f"https://drive.google.com/uc?export=download&id={GDRIVE_FILE_ID}"
    try:
        logger.info("جاري تحميل النموذج من Google Drive...")
        session = requests.Session()
        response = session.get(url, stream=True, timeout=30)

        # التعامل مع صفحة تأكيد الملفات الكبيرة
        for key, value in response.cookies.items():
            if key.startswith("download_warning"):
                confirm_url = f"{url}&confirm={value}"
                response = session.get(confirm_url, stream=True, timeout=30)
                break

        if response.status_code == 200:
            with open(TEMPLATE_PATH, "wb") as f:
                for chunk in response.iter_content(chunk_size=32768):
                    if chunk:
                        f.write(chunk)
            logger.info(f"✅ تم تحميل النموذج بنجاح: {TEMPLATE_PATH}")
            return True
        else:
            logger.error(f"❌ فشل التحميل. كود الحالة: {response.status_code}")
            return False
    except Exception as e:
        logger.error(f"❌ خطأ أثناء تحميل النموذج: {e}")
        return False


# ─── دوال مساعدة ─────────────────────────────────────────────────────────────

def set_rtl_paragraph(paragraph):
    """تفعيل الاتجاه من اليمين لليسار للفقرة"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('{http://schemas.openxmlformats.org/drawingml/2006/main}rtl', '1')


def set_text_rtl(text_frame, text: str, font_size: int = 18, bold: bool = False,
                  color: Optional[tuple] = None, align=PP_ALIGN.RIGHT):
    """كتابة نص عربي مع إعدادات RTL كاملة"""
    text_frame.clear()
    text_frame.word_wrap = True
    para = text_frame.paragraphs[0]
    para.alignment = align
    set_rtl_paragraph(para)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)


def add_bullet_text(text_frame, items: list, font_size: int = 14,
                     title: str = None, title_size: int = 16):
    """إضافة نص على شكل نقاط مع دعم RTL"""
    text_frame.clear()
    text_frame.word_wrap = True

    if title:
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        set_rtl_paragraph(para)
        run = para.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x02, 0x60, 0x99)
    else:
        # إزالة الفقرة الافتراضية الفارغة
        tf_elem = text_frame._txBody
        for p in tf_elem.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
            tf_elem.remove(p)
            break

    for item in items:
        para = text_frame.add_paragraph()
        para.alignment = PP_ALIGN.RIGHT
        set_rtl_paragraph(para)
        run = para.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(font_size)


def parse_multiline(text: str) -> list:
    """تحويل النص المُدخَل إلى قائمة نقاط"""
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    # إذا لم يكن هناك أسطر متعددة، اعتبرها كقائمة بنقطة واحدة
    return lines if lines else [text.strip()]


def build_pptx(data: dict, photos: list) -> io.BytesIO:
    """
    بناء ملف PPTX بناءً على بيانات المستخدم
    data: قاموس يحتوي على بيانات المدرسة والزيارة
    photos: قائمة من bytes لكل صورة
    """
    prs = Presentation(TEMPLATE_PATH)
    slides = prs.slides

    # ── الشريحة 1: صفحة العنوان ──────────────────────────────────────────────
    slide1 = slides[0]
    for shape in slide1.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()

        if "اسم المدرسة" in txt or "..........." in txt:
            set_text_rtl(
                shape.text_frame,
                f"مدرسة {data['school_name']}",
                font_size=24,
                bold=True,
            )
        elif "2025/01/01" in txt or "تاريخ" in txt.lower():
            set_text_rtl(
                shape.text_frame,
                data['visit_date'],
                font_size=16,
            )

    # ── الشريحة 2: ملاحظات الزيارة ──────────────────────────────────────────
    slide2 = slides[1]
    for shape in slide2.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()

        if "ملخص الزيارة" in txt or "ملاحظات تم" in txt:
            # إعادة بناء شريحة الملاحظات بالكامل
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            # ملخص الزيارة
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            set_rtl_paragraph(p)
            run = p.add_run()
            run.text = (
                f"تم زيارة مدرسة {data['school_name']} (الرقم الوزاري: {data['ministry_num']}) "
                f"بتاريخ {data['visit_date']} من قبل فريق شركة تطوير للمباني، "
                f"وتم تدوين الملاحظات الآتية:"
            )
            run.font.size = Pt(13)
            run.font.bold = True

            sections = [
                ("ملاحظات الزيارة:", data['notes']),
                ("الإجراءات المتخذة:", data['actions']),
                ("التوصيات:", data['recommendations']),
            ]
            for section_title, items in sections:
                # عنوان القسم
                sp = tf.add_paragraph()
                sp.alignment = PP_ALIGN.RIGHT
                set_rtl_paragraph(sp)
                sr = sp.add_run()
                sr.text = section_title
                sr.font.size = Pt(14)
                sr.font.bold = True
                sr.font.color.rgb = RGBColor(0x02, 0x60, 0x99)

                for item in items:
                    ip = tf.add_paragraph()
                    ip.alignment = PP_ALIGN.RIGHT
                    set_rtl_paragraph(ip)
                    ir = ip.add_run()
                    ir.text = f"• {item}"
                    ir.font.size = Pt(12)

                # سطر فارغ
                tf.add_paragraph()

            break

    # ── شرائح الصور (كل 3 صور في شريحة) ────────────────────────────────────
    photo_slide_indices = [2, 3]  # الشرائح 3 و4 في النموذج (index 2 و3)

    # تحديد عدد شرائح الصور المطلوبة
    num_photo_slides_needed = (len(photos) + 2) // 3  # كل 3 صور = شريحة

    # إذا كان عدد الشرائح المطلوبة أكثر من المتاح، نضيف شرائح جديدة
    template_photo_slide = slides[photo_slide_indices[0]]  # استخدام أول شريحة صور كقالب

    # إضافة شرائح صور إضافية إذا لزم
    from pptx.util import Inches
    from copy import deepcopy
    import copy

    photo_slides = []

    # استخدام الشرائح الموجودة أولاً
    for idx in photo_slide_indices:
        if idx < len(slides):
            photo_slides.append(slides[idx])

    # إضافة شرائح إضافية إذا لزم
    while len(photo_slides) < num_photo_slides_needed:
        # نسخ شريحة الصور القالب
        xml_src = copy.deepcopy(template_photo_slide._element)
        slide_layout = template_photo_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        # استبدال محتوى الشريحة بنسخة من القالب
        sp_tree = new_slide.shapes._spTree
        for child in list(sp_tree):
            sp_tree.remove(child)
        for child in xml_src.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree/'):
            sp_tree.append(copy.deepcopy(child))
        photo_slides.append(new_slide)

    # تعبئة شرائح الصور
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for slide_idx, photo_slide in enumerate(photo_slides):
        batch_start = slide_idx * 3
        batch = photos[batch_start: batch_start + 3]

        if not batch:
            break

        # تحديث عنوان الشريحة إذا وُجد
        for shape in photo_slide.shapes:
            if shape.has_text_frame and "العنوان" in shape.text_frame.text:
                set_text_rtl(
                    shape.text_frame,
                    f"صور الزيارة - مدرسة {data['school_name']}",
                    font_size=20,
                    bold=True,
                )

        # إزالة أي صور قديمة موجودة في الشريحة (صور placeholder)
        shapes_to_remove = []
        for shape in photo_slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

        # حساب مواضع الصور (3 صور أفقياً)
        margin = Inches(0.3)
        top = Inches(1.5)
        available_w = slide_w - (2 * margin)
        photo_w = (available_w - (2 * Inches(0.15))) // 3
        photo_h = Inches(2.8)

        for i, photo_bytes in enumerate(batch):
            left = margin + i * (photo_w + Inches(0.15))
            img_stream = io.BytesIO(photo_bytes)
            try:
                photo_slide.shapes.add_picture(
                    img_stream, left, top, photo_w, photo_h
                )
            except Exception as e:
                logger.warning(f"تعذّر إضافة الصورة {batch_start + i + 1}: {e}")

    # ── إزالة شرائح الصور غير المستخدمة ─────────────────────────────────────
    # إذا كان عدد الصور أقل من عدد شرائح الصور المتاحة
    used_photo_slides = set(photo_slides[:num_photo_slides_needed])
    slides_to_delete = []
    for idx in photo_slide_indices:
        if idx < len(slides) and slides[idx] not in used_photo_slides:
            slides_to_delete.append(slides[idx])

    for slide in slides_to_delete:
        rId = prs.slides._sldIdLst
        for sldId in rId:
            if prs.slides._sldIdLst_lookup.get(sldId.rId) == slide:
                rId.remove(sldId)
                break

    # ── حفظ الملف في الذاكرة ─────────────────────────────────────────────────
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ─── معالجات المحادثة ─────────────────────────────────────────────────────────

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    """بدء المحادثة"""
    context.user_data.clear()
    await update.message.reply_text(
        "🏫 مرحباً بك في بوت التقارير الفنية\n"
        "━━━━━━━━━━━━━━━━━━━\n"
        "سأساعدك في إنشاء تقرير فني بصيغة PPTX.\n\n"
        "📝 *الخطوة 1/7*\n"
        "أدخل *اسم المدرسة:*",
        parse_mode="Markdown",
    )
    return SCHOOL_NAME


async def get_school_name(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    context.user_data["school_name"] = update.message.text.strip()
    await update.message.reply_text(
        "📝 *الخطوة 2/7*\n"
        "أدخل *الرقم الوزاري* للمدرسة:",
        parse_mode="Markdown",
    )
    return MINISTRY_NUM


async def get_ministry_num(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    context.user_data["ministry_num"] = update.message.text.strip()
    await update.message.reply_text(
        "📝 *الخطوة 3/7*\n"
        "أدخل *تاريخ الزيارة* (مثال: 15/03/2025):",
        parse_mode="Markdown",
    )
    return VISIT_DATE


async def get_visit_date(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    context.user_data["visit_date"] = update.message.text.strip()
    await update.message.reply_text(
        "📝 *الخطوة 4/7*\n"
        "أدخل *الملاحظات* (كل ملاحظة في سطر جديد):",
        parse_mode="Markdown",
    )
    return NOTES


async def get_notes(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    context.user_data["notes"] = parse_multiline(update.message.text)
    await update.message.reply_text(
        "📝 *الخطوة 5/7*\n"
        "أدخل *الإجراءات المتخذة* (كل إجراء في سطر جديد):",
        parse_mode="Markdown",
    )
    return ACTIONS


async def get_actions(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    context.user_data["actions"] = parse_multiline(update.message.text)
    await update.message.reply_text(
        "📝 *الخطوة 6/7*\n"
        "أدخل *التوصيات* (كل توصية في سطر جديد):",
        parse_mode="Markdown",
    )
    return RECOMMENDATIONS


async def get_recommendations(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    context.user_data["recommendations"] = parse_multiline(update.message.text)
    context.user_data["photos"] = []
    await update.message.reply_text(
        "📝 *الخطوة 7/7 - الصور*\n"
        "أرسل صور الزيارة واحدة تلو الأخرى.\n"
        "━━━━━━━━━━━━━━━━━━━\n"
        "⚠️ سيتم ترتيب كل 3 صور في شريحة واحدة.\n\n"
        "عند الانتهاء من إرسال الصور، أرسل:\n"
        "✅ /done — لإنشاء التقرير\n"
        "❌ /cancel — للإلغاء",
        parse_mode="Markdown",
    )
    return PHOTOS


async def get_photo(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    """استقبال صورة وتخزينها"""
    photos = context.user_data.setdefault("photos", [])

    # أخذ أعلى جودة متاحة
    photo_file = await update.message.photo[-1].get_file()
    photo_bytes = await photo_file.download_as_bytearray()
    photos.append(bytes(photo_bytes))

    count = len(photos)
    slides_count = (count + 2) // 3
    await update.message.reply_text(
        f"✅ تم استلام الصورة *{count}*\n"
        f"📊 ستُوزَّع على *{slides_count}* شريحة\n\n"
        "أرسل صورة أخرى أو /done للإنهاء.",
        parse_mode="Markdown",
    )
    return PHOTOS


async def done_photos(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    """إنهاء استقبال الصور وإنشاء التقرير"""
    data = context.user_data
    photos = data.get("photos", [])

    school = data.get("school_name", "—")
    ministry = data.get("ministry_num", "—")
    visit = data.get("visit_date", "—")
    notes_count = len(data.get("notes", []))
    actions_count = len(data.get("actions", []))
    recs_count = len(data.get("recommendations", []))

    await update.message.reply_text(
        f"📋 *ملخص التقرير*\n"
        f"━━━━━━━━━━━━━━━━━━━\n"
        f"🏫 المدرسة: {school}\n"
        f"🔢 الرقم الوزاري: {ministry}\n"
        f"📅 تاريخ الزيارة: {visit}\n"
        f"📌 ملاحظات: {notes_count} بند\n"
        f"⚙️ إجراءات: {actions_count} بند\n"
        f"💡 توصيات: {recs_count} بند\n"
        f"🖼️ صور: {len(photos)} صورة\n"
        f"━━━━━━━━━━━━━━━━━━━\n"
        f"⏳ جاري إنشاء التقرير...",
        parse_mode="Markdown",
    )

    try:
        pptx_bytes = build_pptx(data, photos)
        filename = f"تقرير_{school}_{visit.replace('/', '-')}.pptx"
        await update.message.reply_document(
            document=pptx_bytes,
            filename=filename,
            caption=(
                f"✅ *تم إنشاء التقرير الفني بنجاح*\n"
                f"🏫 مدرسة: {school}\n"
                f"📅 تاريخ: {visit}"
            ),
            parse_mode="Markdown",
        )
    except Exception as e:
        logger.error(f"خطأ في إنشاء التقرير: {e}", exc_info=True)
        await update.message.reply_text(
            f"❌ حدث خطأ أثناء إنشاء التقرير:\n`{e}`\n\n"
            "تأكد من أن النموذج محمّل بشكل صحيح وأعد المحاولة.",
            parse_mode="Markdown",
        )

    return ConversationHandler.END


async def cancel(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    """إلغاء العملية"""
    context.user_data.clear()
    await update.message.reply_text(
        "❌ تم إلغاء العملية.\n"
        "أرسل /start للبدء من جديد.",
        reply_markup=ReplyKeyboardRemove(),
    )
    return ConversationHandler.END


async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "📖 *مساعدة البوت*\n"
        "━━━━━━━━━━━━━━━━━━━\n"
        "/start — بدء إنشاء تقرير جديد\n"
        "/done  — إنهاء إرسال الصور وإنشاء التقرير\n"
        "/cancel — إلغاء العملية الحالية\n"
        "/help  — عرض هذه المساعدة\n"
        "━━━━━━━━━━━━━━━━━━━\n"
        "💡 *ملاحظات:*\n"
        "• أرسل الملاحظات والإجراءات والتوصيات كل بند في سطر منفصل\n"
        "• يمكنك إرسال أي عدد من الصور (كل 3 في شريحة)\n"
        "• الصور اختيارية، يمكنك إرسال /done مباشرة بدون صور",
        parse_mode="Markdown",
    )


# ─── نقطة الدخول الرئيسية ────────────────────────────────────────────────────

def main() -> None:
    # تحميل النموذج عند التشغيل
    if not os.path.exists(TEMPLATE_PATH):
        if not download_template():
            logger.warning(
                "⚠️  لم يتم تحميل النموذج. تأكد من صلاحيات الملف على Google Drive "
                "(يجب أن يكون عاماً أو قابلاً للتنزيل)"
            )
    else:
        logger.info(f"✅ النموذج موجود مسبقاً: {TEMPLATE_PATH}")

    app = Application.builder().token(TOKEN).build()

    conv_handler = ConversationHandler(
        entry_points=[CommandHandler("start", start)],
        states={
            SCHOOL_NAME:      [MessageHandler(filters.TEXT & ~filters.COMMAND, get_school_name)],
            MINISTRY_NUM:     [MessageHandler(filters.TEXT & ~filters.COMMAND, get_ministry_num)],
            VISIT_DATE:       [MessageHandler(filters.TEXT & ~filters.COMMAND, get_visit_date)],
            NOTES:            [MessageHandler(filters.TEXT & ~filters.COMMAND, get_notes)],
            ACTIONS:          [MessageHandler(filters.TEXT & ~filters.COMMAND, get_actions)],
            RECOMMENDATIONS:  [MessageHandler(filters.TEXT & ~filters.COMMAND, get_recommendations)],
            PHOTOS: [
                MessageHandler(filters.PHOTO, get_photo),
                CommandHandler("done", done_photos),
            ],
        },
        fallbacks=[CommandHandler("cancel", cancel)],
        allow_reentry=True,
    )

    app.add_handler(conv_handler)
    app.add_handler(CommandHandler("help", help_command))
    app.add_handler(CommandHandler("cancel", cancel))

    logger.info("🤖 البوت يعمل الآن...")
    app.run_polling(allowed_updates=Update.ALL_TYPES)


if __name__ == "__main__":
    main()
